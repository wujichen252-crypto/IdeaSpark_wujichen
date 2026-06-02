"""File-related business logic — OSS upload + export."""
import logging

from apps.files import oss

logger = logging.getLogger(__name__)


def upload(file_obj) -> dict:
    """Upload file to OSS."""
    filename = getattr(file_obj, 'name', None)
    return oss.upload_file(file_obj, filename)


def export_docx(html_content: str, file_name: str = 'document.docx') -> bytes:
    """Export HTML content to DOCX."""
    from docx import Document
    from docx.shared import Inches, Pt
    from docx.oxml.ns import qn
    import re

    doc = Document()

    # Set default font
    style = doc.styles['Normal']
    font = style.font
    font.name = 'Arial'
    style.element.rPr.rFonts.set(qn('w:eastAsia'), '微软雅黑')

    # Simple HTML to DOCX conversion
    lines = html_content.replace('<br>', '\n').replace('<br/>', '\n').replace('<br />', '\n')
    # Remove HTML tags
    text = re.sub(r'<[^>]+>', '', lines)

    for paragraph_text in text.split('\n'):
        paragraph_text = paragraph_text.strip()
        if not paragraph_text:
            continue

        # Detect headers
        if paragraph_text.startswith('# '):
            doc.add_heading(paragraph_text[2:], level=1)
        elif paragraph_text.startswith('## '):
            doc.add_heading(paragraph_text[3:], level=2)
        elif paragraph_text.startswith('### '):
            doc.add_heading(paragraph_text[4:], level=3)
        else:
            doc.add_paragraph(paragraph_text)

    import io
    buf = io.BytesIO()
    doc.save(buf)
    buf.seek(0)
    return buf.getvalue()


def export_xlsx(data: list, file_name: str = 'workbook.xlsx', sheet_name: str = 'Sheet1') -> bytes:
    """Export data to XLSX."""
    from openpyxl import Workbook

    wb = Workbook()
    ws = wb.active
    ws.title = sheet_name

    if data:
        for row_idx, row_data in enumerate(data, 1):
            for col_idx, cell_value in enumerate(row_data, 1):
                ws.cell(row=row_idx, column=col_idx, value=cell_value)

    import io
    buf = io.BytesIO()
    wb.save(buf)
    buf.seek(0)
    return buf.getvalue()


def export_pptx(slides: list, file_name: str = 'presentation.pptx') -> bytes:
    """Export slides data to PPTX."""
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()

    for slide_data in slides:
        slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title and Content

        title = slide_data.get('title', '') if isinstance(slide_data, dict) else ''
        content = slide_data.get('content', '') if isinstance(slide_data, dict) else str(slide_data)

        if hasattr(slide, 'shapes'):
            if title and slide.shapes.title:
                slide.shapes.title.text = title
            if content and slide.placeholders:
                try:
                    placeholder = slide.placeholders[1]
                    placeholder.text = content
                except (IndexError, AttributeError):
                    pass

    import io
    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.getvalue()
